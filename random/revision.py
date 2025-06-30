from pptx import Presentation
from pptx.util import Inches

# Create a presentation
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]

# Title slide
slide_0 = prs.slides.add_slide(title_slide_layout)
slide_0.shapes.title.text = "Chapter 8: Functions in Python"
slide_0.placeholders[1].text = "With a Focus on Recursion"

# Slide 1: What is a Function?
slide_1 = prs.slides.add_slide(content_slide_layout)
slide_1.shapes.title.text = "What is a Function?"
slide_1.placeholders[1].text = (
    "Functions allow reuse of code blocks.\n"
    "Avoids repeating the same logic multiple times.\n"
    "Syntax:\n"
    "def greet():\n"
    "    print('Hello')"
)

# Slide 2: Function Types and Arguments
slide_2 = prs.slides.add_slide(content_slide_layout)
slide_2.shapes.title.text = "Function Types & Arguments"
slide_2.placeholders[1].text = (
    "• Built-in functions: pre-defined in Python.\n"
    "• User-defined functions: defined by users.\n\n"
    "Function with arguments: allows passing inputs.\n"
    "Default parameters: used if no input is passed."
)

# Slide 3: Use of return Statement
slide_3 = prs.slides.add_slide(content_slide_layout)
slide_3.shapes.title.text = "Use of return Statement"
slide_3.placeholders[1].text = (
    "• return sends the output back to the caller.\n"
    "• Useful when output needs to be reused or stored.\n"
    "• print() only displays the result.\n\n"
    "Example:\n"
    "def square(x):\n"
    "    return x * x"
)

# Slide 4: Introduction to Recursion
slide_4 = prs.slides.add_slide(content_slide_layout)
slide_4.shapes.title.text = "What is Recursion?"
slide_4.placeholders[1].text = (
    "Recursion is a function calling itself.\n"
    "Used to solve problems using smaller instances.\n"
    "Important to define a base case to stop recursion.\n\n"
    "Common in factorial, Fibonacci, tree traversal, etc."
)

# Slide 5: Example – Factorial Using Recursion
slide_5 = prs.slides.add_slide(content_slide_layout)
slide_5.shapes.title.text = "Factorial Function (Recursive)"
slide_5.placeholders[1].text = (
    "def factorial(n):\n"
    "    if n == 0 or n == 1:\n"
    "        return 1  # base case\n"
    "    else:\n"
    "        return n * factorial(n - 1)  # recursive call\n\n"
    "Calling factorial(5) returns 5 * 4 * 3 * 2 * 1 = 120"
)

# Slide 6: How Recursion Works (Step-by-step)
slide_6 = prs.slides.add_slide(content_slide_layout)
slide_6.shapes.title.text = "How Recursion Works"
slide_6.placeholders[1].text = (
    "factorial(3)\n"
    "→ 3 * factorial(2)\n"
    "→ 3 * 2 * factorial(1)\n"
    "→ 3 * 2 * 1 (base case returns 1)\n"
    "= 6\n\n"
    "Each call is paused until the inner call finishes."
)

# Slide 7: Caution with Recursion
slide_7 = prs.slides.add_slide(content_slide_layout)
slide_7.shapes.title.text = "Caution with Recursion"
slide_7.placeholders[1].text = (
    "• Always define a base case to avoid infinite loops.\n"
    "• Recursive functions can cause stack overflow if too deep.\n"
    "• Prefer iteration if recursion is not required.\n\n"
    "Use recursion only when the problem is naturally recursive."
)

# Save the presentation
pptx_path = "/mnt/data/Chapter_8_Functions_with_Recursion.pptx"
prs.save(pptx_path)

pptx_path
